import io
import math
import numpy as np
from typing import List, Dict, Any, Tuple
import google.generativeai as genai
import openai
from pypdf import PdfReader
import docx
from pptx import Presentation
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer

# Recursive Chunker
def chunk_text(text: str, page_number: int = 1, chunk_size: int = 800, overlap: int = 100) -> List[Dict[str, Any]]:
    """
    Splits text into chunks of roughly chunk_size characters with overlap.
    Returns a list of dicts: {'text': chunk_text, 'page': page_number}
    """
    chunks = []
    text_len = len(text)
    if text_len <= chunk_size:
        if text.strip():
            chunks.append({"text": text.strip(), "page": page_number})
        return chunks
        
    start = 0
    while start < text_len:
        end = start + chunk_size
        
        # Try to find a nice boundary (newline or space) near the end
        if end < text_len:
            # Look back up to 'overlap' characters for a sentence boundary or newline
            boundary = text.rfind('\n', end - overlap, end)
            if boundary == -1:
                boundary = text.rfind('. ', end - overlap, end)
            if boundary == -1:
                boundary = text.rfind(' ', end - overlap, end)
            
            if boundary != -1:
                end = boundary + 1
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append({"text": chunk, "page": page_number})
            
        start = end - overlap
        if start >= text_len or end >= text_len:
            break
            
    return chunks

class RAGService:
    @staticmethod
    def parse_document(file_content: bytes, file_type: str) -> List[Dict[str, Any]]:
        """
        Parses document bytes based on file extension.
        Returns a list of chunks: [{'text': str, 'page': int}]
        """
        chunks = []
        file_like = io.BytesIO(file_content)
        
        try:
            if file_type == "pdf":
                reader = PdfReader(file_like)
                for idx, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        chunks.extend(chunk_text(page_text, page_number=idx + 1))
                        
            elif file_type == "docx":
                doc = docx.Document(file_like)
                full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                chunks.extend(chunk_text(full_text, page_number=1))
                
            elif file_type == "pptx":
                prs = Presentation(file_like)
                for idx, slide in enumerate(prs.slides):
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text + "\n"
                    if slide_text.strip():
                        chunks.extend(chunk_text(slide_text, page_number=idx + 1))
                        
            elif file_type in ["xlsx", "xls"]:
                # Convert Excel sheets to text descriptions
                xls = pd.ExcelFile(file_like)
                for sheet_name in xls.sheet_names:
                    df = xls.parse(sheet_name)
                    # Convert to csv-like representation
                    csv_text = df.to_csv(index=False)
                    desc = f"Sheet: {sheet_name}\n" + csv_text
                    chunks.extend(chunk_text(desc, page_number=1, chunk_size=1000))
                    
            elif file_type == "csv":
                df = pd.read_csv(file_like)
                csv_text = df.to_csv(index=False)
                chunks.extend(chunk_text(csv_text, page_number=1, chunk_size=1000))
                
            else: # txt and general fallback
                text_content = file_content.decode("utf-8", errors="ignore")
                chunks.extend(chunk_text(text_content, page_number=1))
                
        except Exception as e:
            print(f"Error parsing document type {file_type}: {e}")
            # Try plain text fallback
            try:
                text_content = file_content.decode("utf-8", errors="ignore")
                chunks.extend(chunk_text(text_content, page_number=1))
            except:
                pass
                
        return chunks

    @staticmethod
    def generate_embeddings(texts: List[str], api_settings: Dict[str, str]) -> List[List[float]]:
        """
        Generates vector embeddings for a list of texts.
        Supports:
          - Gemini (if API Key provided)
          - OpenAI (if API Key provided)
          - TF-IDF local vectorizer (as lightweight fallback)
        """
        gemini_key = api_settings.get("GEMINI_API_KEY")
        openai_key = api_settings.get("OPENAI_API_KEY")
        
        # 1. Gemini
        if gemini_key:
            try:
                genai.configure(api_key=gemini_key)
                response = genai.embed_content(
                    model="models/text-embedding-004",
                    content=texts,
                    task_type="retrieval_document"
                )
                return response['embedding']
            except Exception as e:
                print(f"Gemini embedding error, falling back: {e}")
                
        # 2. OpenAI
        if openai_key:
            try:
                client = openai.OpenAI(api_key=openai_key)
                response = client.embeddings.create(
                    input=texts,
                    model="text-embedding-3-small"
                )
                return [x.embedding for x in response.data]
            except Exception as e:
                print(f"OpenAI embedding error, falling back: {e}")
                
        # 3. TF-IDF fallback: Generate pseudo-embeddings (size 128 for compatibility)
        # We compute a local term frequency vector and pad/hash it to 128 dimensions.
        return RAGService._generate_local_embeddings(texts)

    @staticmethod
    def _generate_local_embeddings(texts: List[str]) -> List[List[float]]:
        """
        Generates a 128-dimensional mock embedding using TF-IDF + hashing.
        Allows offline cosine-similarity calculation that actually reflects text overlap!
        """
        if not texts:
            return []
            
        # Add basic vocabulary helper
        vectorizer = TfidfVectorizer(max_features=128, stop_words='english')
        try:
            # Fit/transform
            tfidf_matrix = vectorizer.fit_transform(texts).toarray()
            # If TF-IDF returns fewer than 128 features, pad with zeros
            features_found = tfidf_matrix.shape[1]
            if features_found < 128:
                padding = np.zeros((len(texts), 128 - features_found))
                tfidf_matrix = np.hstack((tfidf_matrix, padding))
            # Normalize each vector
            norms = np.linalg.norm(tfidf_matrix, axis=1, keepdims=True)
            norms[norms == 0] = 1.0  # avoid division by zero
            normalized_embeddings = (tfidf_matrix / norms).tolist()
            return normalized_embeddings
        except Exception as e:
            # Fallback to simple hash-based pseudo-embeddings if fitting fails
            print(f"Local TF-IDF fitting failed, falling back to hash embeddings: {e}")
            embeddings = []
            for text in texts:
                vec = [0.0] * 128
                words = text.lower().split()
                for word in words:
                    h = hash(word) % 128
                    vec[h] += 1.0
                # L2 Normalize
                vec_sum = sum(x*x for x in vec)
                norm = math.sqrt(vec_sum) if vec_sum > 0 else 1.0
                embeddings.append([x / norm for x in vec])
            return embeddings

    @staticmethod
    def retrieve_chunks(
        query: str, 
        all_chunks: List[Dict[str, Any]], 
        user_role: str, 
        api_settings: Dict[str, str],
        top_k: int = 5
    ) -> List[Dict[str, Any]]:
        """
        Retrieves top_k chunks relevant to query, checking role permissions:
        - Admin: access to all chunks
        - HR: access to 'General' and 'HR' chunks
        - Finance: access to 'General' and 'Finance' chunks
        - General: access to 'General' chunks only
        """
        # 1. Filter by Role Permissions
        allowed_roles = ["General"]
        if user_role == "Admin":
            allowed_roles = ["General", "HR", "Finance", "Admin"]
        elif user_role == "HR":
            allowed_roles = ["General", "HR"]
        elif user_role == "Finance":
            allowed_roles = ["General", "Finance"]
            
        filtered_chunks = [
            c for c in all_chunks 
            if c.get("access_role", "General") in allowed_roles
        ]
        
        if not filtered_chunks:
            return []
            
        # 2. Get Query Embedding
        query_emb = RAGService.generate_embeddings([query], api_settings)
        if not query_emb:
            # Fallback to simple text search if embedding generation failed
            return RAGService._keyword_retrieve(query, filtered_chunks, top_k)
            
        q_vec = np.array(query_emb[0])
        
        # 3. Calculate Cosine Similarity
        scored_chunks = []
        for chunk in filtered_chunks:
            c_vec_raw = chunk.get("embedding")
            if c_vec_raw:
                c_vec = np.array(c_vec_raw)
                # Compute cosine similarity
                dot_product = np.dot(q_vec, c_vec)
                q_norm = np.linalg.norm(q_vec)
                c_norm = np.linalg.norm(c_vec)
                similarity = dot_product / (q_norm * c_norm) if (q_norm * c_norm) > 0 else 0
            else:
                similarity = 0.0
            
            scored_chunks.append((similarity, chunk))
            
        # Sort by similarity score descending
        scored_chunks.sort(key=lambda x: x[0], reverse=True)
        
        results = []
        for score, chunk in scored_chunks[:top_k]:
            chunk_copy = chunk.copy()
            chunk_copy["score"] = float(score)
            results.append(chunk_copy)
            
        return results

    @staticmethod
    def _keyword_retrieve(query: str, chunks: List[Dict[str, Any]], top_k: int) -> List[Dict[str, Any]]:
        # Sub-fallback: string keyword matching
        words = set(query.lower().split())
        scored = []
        for chunk in chunks:
            text = chunk.get("text", "").lower()
            score = sum(1.0 for w in words if w in text)
            scored.append((score, chunk))
        scored.sort(key=lambda x: x[0], reverse=True)
        return [c[1] for c in scored[:top_k] if c[0] > 0]

    @staticmethod
    def generate_response(
        query: str, 
        retrieved_chunks: List[Dict[str, Any]], 
        api_settings: Dict[str, str]
    ) -> Tuple[str, List[Dict[str, Any]]]:
        """
        Generates the LLM response given a query and retrieved chunks context.
        Provides citations back to original files.
        """
        gemini_key = api_settings.get("GEMINI_API_KEY")
        openai_key = api_settings.get("OPENAI_API_KEY")
        
        # Build prompt context
        context_str = ""
        citations = []
        for idx, chunk in enumerate(retrieved_chunks):
            doc_name = chunk.get("file_name", "Unknown File")
            page_num = chunk.get("page_number", 1)
            context_str += f"[Source {idx+1}]: {doc_name} (Page {page_num})\nContent: {chunk.get('text', '')}\n\n"
            citations.append({
                "source_id": idx + 1,
                "file_name": doc_name,
                "page_number": page_num,
                "snippet": chunk.get('text', '')[:200] + "..."
            })
            
        system_prompt = (
            "You are NeuroRAG, an intelligent corporate knowledge assistant.\n"
            "Use ONLY the following context sources to answer the query. "
            "If the answer cannot be found in the context, state that you do not know. "
            "Cite your sources in the text using bracketed numbers like [1] or [2] matching the source index.\n\n"
            f"Context Sources:\n{context_str}"
        )
        
        user_prompt = f"User Query: {query}"
        
        # 1. Real LLM Call: Gemini
        if gemini_key:
            try:
                genai.configure(api_key=gemini_key)
                model = genai.GenerativeModel("gemini-1.5-flash")
                # Combine prompts
                full_prompt = f"{system_prompt}\n\n{user_prompt}"
                response = model.generate_content(full_prompt)
                return response.text, citations
            except Exception as e:
                print(f"Gemini generation error: {e}")
                
        # 2. Real LLM Call: OpenAI
        if openai_key:
            try:
                client = openai.OpenAI(api_key=openai_key)
                response = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ]
                )
                return response.choices[0].message.content, citations
            except Exception as e:
                print(f"OpenAI generation error: {e}")
                
        # 3. Smart Demo Response Fallback (when no keys provided)
        # Parse context to synthesize a realistic reply
        reply = RAGService._synthesize_demo_response(query, retrieved_chunks)
        return reply, citations

    @staticmethod
    def _synthesize_demo_response(query: str, retrieved_chunks: List[Dict[str, Any]]) -> str:
        """
        Synthesizes a clean demo reply citing context if available.
        """
        import re

        # Normalize and clean query
        clean_query = query.strip().lower().rstrip('?.!')
        
        # 1. Handle Greetings / Chit-Chat
        greetings = {"hi", "hello", "hey", "greetings", "good morning", "good afternoon", "good evening", "howdy", "sup", "hola"}
        if clean_query in greetings:
            return (
                "Hello! I am NeuroRAG, your intelligent corporate knowledge assistant. "
                "I am ready to search and answer questions based on your uploaded company documents. "
                "How can I help you today?"
            )
            
        if any(phrase in clean_query for phrase in ["who are you", "what is your name", "your name"]):
            return (
                "I am NeuroRAG, an AI-powered Enterprise Knowledge Assistant. "
                "I can help you search, summarize, and query company documents with role-based security access control. "
                "How can I assist you today?"
            )
            
        if not retrieved_chunks:
            return (
                "Based on my search of the uploaded company documents, I couldn't find any information "
                "regarding your query. Please make sure the relevant document is uploaded and that you have "
                "the necessary security clearance to access it."
            )

        # 2. Extract keywords from the query
        STOPWORDS = {
            'a', 'about', 'above', 'after', 'again', 'against', 'all', 'am', 'an', 'and', 'any', 'are', 'arent', 'as', 'at',
            'be', 'because', 'been', 'before', 'being', 'below', 'between', 'both', 'but', 'by', 'cant', 'cannot', 'could',
            'couldnt', 'did', 'didnt', 'do', 'does', 'doesnt', 'doing', 'dont', 'down', 'during', 'each', 'few', 'for', 'from',
            'further', 'had', 'hadnt', 'has', 'hasnt', 'have', 'havent', 'having', 'he', 'hed', 'hell', 'hes', 'her', 'here',
            'heres', 'hers', 'herself', 'him', 'himself', 'his', 'how', 'hows', 'i', 'id', 'ill', 'im', 'ive', 'if', 'in',
            'into', 'is', 'isnt', 'it', 'its', 'itself', 'lets', 'me', 'more', 'most', 'mustnt', 'my', 'myself', 'no', 'nor',
            'not', 'of', 'off', 'on', 'once', 'only', 'or', 'other', 'ought', 'our', 'ours', 'ourselves', 'out', 'over', 'own',
            'same', 'shant', 'she', 'shed', 'shell', 'shes', 'should', 'shouldnt', 'so', 'some', 'such', 'than', 'that',
            'thats', 'the', 'their', 'theirs', 'them', 'themselves', 'then', 'there', 'theres', 'these', 'they', 'theyd',
            'theyll', 'theyre', 'theyve', 'this', 'those', 'through', 'to', 'too', 'under', 'until', 'up', 'very', 'was',
            'wasnt', 'we', 'wed', 'well', 'were', 'weve', 'werent', 'what', 'whats', 'when', 'whens', 'where', 'wheres',
            'which', 'while', 'who', 'whos', 'whom', 'why', 'whys', 'with', 'wont', 'would', 'wouldnt', 'you', 'youd',
            'youll', 'youre', 'youve', 'your', 'yours', 'yourself', 'yourselves',
            'provide', 'show', 'tell', 'get', 'give', 'list', 'please', 'some', 'any', 'ask', 'asking', 'think'
        }
        
        query_clean = re.sub(r'[^\w\s]', ' ', query.lower())
        query_words = [w for w in query_clean.split() if w and w not in STOPWORDS]
        
        stemmed_words = set()
        for w in query_words:
            stemmed_words.add(w)
            if w.endswith('s') and len(w) > 3:
                stemmed_words.add(w[:-1])
            if w.endswith('es') and len(w) > 4:
                stemmed_words.add(w[:-2])
            if 'resume' in w:
                stemmed_words.add('resume')

        # 3. Extract and score segments from all retrieved chunks
        scored_segments = []
        is_question_query = any(qw in clean_query for qw in ["question", "interview", "query", "ask", "solve", "problem"])
        
        for chunk_idx, chunk in enumerate(retrieved_chunks):
            text = chunk.get("text", "")
            doc_name = chunk.get("file_name", "document")
            page = chunk.get("page_number", 1)
            
            # Split text into lines/sentences
            raw_lines = text.split('\n')
            segments = []
            for line in raw_lines:
                line = line.strip()
                if not line:
                    continue
                if len(line) < 150:
                    segments.append(line)
                else:
                    sub_segs = re.split(r'(?<=[.!?])\s+', line)
                    for s in sub_segs:
                        s_clean = s.strip()
                        if s_clean:
                            segments.append(s_clean)
                            
            for seg in segments:
                if len(seg) < 10:
                    continue
                # Score segment based on stemmed word overlap
                overlap_score = 0
                seg_lower = seg.lower()
                seg_words = set(re.sub(r'[^\w\s]', ' ', seg_lower).split())
                
                # Check for word presence and full word match
                for word in stemmed_words:
                    if word in seg_lower:
                        overlap_score += 2
                    if word in seg_words:
                        overlap_score += 3
                        
                # Extra points if query is seeking questions and segment looks like a question or list item
                if is_question_query:
                    if "?" in seg:
                        overlap_score += 8
                    elif re.match(r'^\d+[\.\)]', seg):
                        overlap_score += 6
                    elif seg_lower.startswith(("what ", "how ", "why ", "write ", "explain ", "describe ")):
                        overlap_score += 6
                        
                if seg.startswith(("•", "-", "*")):
                    overlap_score += 1
                    
                if overlap_score > 0:
                    scored_segments.append({
                        "score": overlap_score,
                        "text": seg,
                        "doc_name": doc_name,
                        "page": page,
                        "chunk_idx": chunk_idx
                    })
                    
        # Sort by score descending
        scored_segments.sort(key=lambda x: x["score"], reverse=True)
        
        # 4. Synthesize reply from top scored segments
        if scored_segments and scored_segments[0]["score"] >= 3:
            # Group top segments (up to 6) by source document to present nicely
            selected_segs = []
            seen_texts = set()
            for seg in scored_segments:
                # Deduplicate very similar texts
                norm_text = re.sub(r'\s+', '', seg["text"].lower())
                if norm_text not in seen_texts:
                    seen_texts.add(norm_text)
                    selected_segs.append(seg)
                    if len(selected_segs) >= 6:
                        break
                        
            # Group by document
            doc_groups = {}
            for seg in selected_segs:
                doc_name = seg["doc_name"]
                if doc_name not in doc_groups:
                    doc_groups[doc_name] = []
                doc_groups[doc_name].append(seg)
                
            response = "Based on the retrieved company documentation, here are the relevant details:\n\n"
            for doc_name, segs in doc_groups.items():
                response += f"From **{doc_name}**:\n"
                for seg in segs:
                    page_str = f" (Page {seg['page']})" if seg['page'] else ""
                    # If it doesn't already start with a list marker, add a bullet point
                    bullet = "" if seg['text'].startswith(("•", "-", "*", "1.", "2.", "3.", "4.", "5.", "6.", "7.", "8.", "9.", "10.")) else "- "
                    response += f"  {bullet}{seg['text']}{page_str} [{seg['chunk_idx'] + 1}]\n"
                response += "\n"
                
            response += "Is there anything specific you would like to clarify or explore further regarding these details?"
            return response

        # 5. Default Summary Fallback if no specific segments score highly
        doc_names = list(set([c.get("file_name") for c in retrieved_chunks]))
        response_template = (
            f"Based on the company documentation retrieved (specifically from {', '.join(doc_names)}):\n\n"
        )
        for idx, chunk in enumerate(retrieved_chunks[:3]):
            text = chunk.get("text", "")
            doc_name = chunk.get("file_name", "document")
            page = chunk.get("page_number", 1)
            # Find first 2 lines or sentences
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            summary_lines = lines[:2]
            summary_text = " | ".join(summary_lines)
            if len(summary_text) > 200:
                summary_text = summary_text[:200] + "..."
            response_template += f"- According to **{doc_name}** (Page {page}): {summary_text} [{idx+1}]\n"
            
        response_template += (
            f"\nIs there anything specific you would like to explore regarding these details?"
        )
        return response_template
